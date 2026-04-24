"""Optional PPTX exporter. Requires `python-pptx` (install via [pptx] extra)."""
from __future__ import annotations

from pathlib import Path

from ..models.domain import AnalysisReport


def write_pptx_report(report: AnalysisReport, out_path: Path) -> Path:
    """Render a summary deck. Raises ImportError if python-pptx is not installed."""
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:  # pragma: no cover
        raise ImportError(
            "python-pptx is required for PPTX export. Install with: pip install model-upgrade-analyzer[pptx]"
        ) from exc

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Model Upgrade Analyzer"
    title_slide.placeholders[1].text = f"Generated {report.generated_at}\n{report.repo_path}"

    # Summary slide
    summary = prs.slides.add_slide(prs.slide_layouts[1])
    summary.shapes.title.text = "Summary"
    body = summary.placeholders[1].text_frame
    body.text = f"Deployments: {len(report.deployments)}"
    for label, count in [
        ("Model cards", len(report.model_cards)),
        ("Code references", len(report.code_references)),
        ("Prompt observations", len(report.prompt_observations)),
        ("Global findings", len(report.findings)),
    ]:
        p = body.add_paragraph()
        p.text = f"{label}: {count}"

    # Per-deployment slides (cap to avoid gigantic decks)
    for d in report.deployments[:50]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = d.deployment_name
        tf = slide.placeholders[1].text_frame
        tf.text = f"Current: {d.current_model or 'unknown'}"
        for line in [
            f"Target: {d.target_model or 'unknown'}",
            f"Urgency: {d.urgency.value}",
            f"Complexity: {d.migration_complexity}",
            f"Code refs: {len(d.code_references)}  |  Prompts: {len(d.prompt_files)}",
        ]:
            p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(14)

    prs.save(str(out_path))
    return out_path
